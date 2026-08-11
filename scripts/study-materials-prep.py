#!/usr/bin/env python3
"""
study-materials-prep.py — Prepare study materials for RAG indexing.

Recursively scans a source directory, extracts content from all supported
files (archives, PDFs, DOCX, PPTX, TXT, MD, FB2, DJVU, HTML, images),
and produces structured Markdown output ready for AI agents.

Within this repository, these prepared Markdown outputs are shared upstream
academic infrastructure for:
- lecture-transcript
- homework-management
- case-analyzer

The launching agent/workflow must verify the resulting pack before trusting it
and keep originals available as fallback when conversion quality is weak.

Usage:
    python study-materials-prep.py --source <path>
    python study-materials-prep.py --source <path> --output <path>
    python study-materials-prep.py --source <path> --install-deps

The subject name is auto-detected from the source directory name.
Output goes to ./study-output/<subject>/ by default (or --output if given).
OCR is always enabled for PDFs and images.
Use --install-deps to auto-install missing dependencies.
"""

import argparse
from collections import defaultdict
import hashlib
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from datetime import datetime, timezone
from pathlib import Path

# ---------------------------------------------------------------------------
LOG_FORMAT = "[%(levelname)s] %(message)s"
logging.basicConfig(level=logging.INFO, format=LOG_FORMAT)
log = logging.getLogger("study-prep")

ARCHIVE_EXTS = {".zip", ".rar", ".7z", ".tar", ".tar.gz", ".tgz"}
TEXT_EXTS = {".txt", ".md", ".markdown", ".text", ".log", ".csv"}
DOC_EXTS = {".docx", ".doc"}
PDF_EXTS = {".pdf"}
PPTX_EXTS = {".pptx", ".ppt"}
BOOK_EXTS = {".epub", ".fb2", ".mobi"}
HTML_EXTS = {".html", ".htm", ".xhtml"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
SKIP_EXTS = {".url", ".lnk", ".ini", ".tmp", ".DS_Store", ".Thumbs.db",
             ".json", ".jsonl", ".orig", ".bak", ".py", ".pyc", ".exe",
             ".dll", ".so", ".dylib"}
SKIP_DIRS = {".git", ".qwen", ".claude", ".codex", ".cursor", ".gemini",
             ".opencode", ".scratchpad", "__pycache__", "node_modules",
             ".agent-memory", "coordination", "knowledge-base",
             "_codex_tmp", "chrome", "chromedriver", "geckodriver",
             "tools", "ДР УП-52 2026", "ДЗ_материалы"}


def check_deps():
    """Check optional dependencies, return warnings list."""
    w = []
    for pkg, pip, desc in [
        ("docx", "python-docx", "DOCX extraction"),
        ("pptx", "python-pptx", "PPTX extraction"),
        ("PIL", "Pillow", "Image/OCR support"),
    ]:
        try:
            __import__(pkg)
        except ImportError:
            w.append(f"MISSING: {pip} ({desc}) — pip install {pip}")
    if shutil.which("tesseract"):
        try:
            import pytesseract  # noqa
        except ImportError:
            w.append("MISSING: pytesseract (OCR) — pip install pytesseract")
    else:
        w.append("OPTIONAL: tesseract-ocr not on PATH (OCR disabled)")
    return w


# ---------------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------------

def read_text(path: Path) -> str:
    for enc in ("utf-8", "cp1251", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, ValueError, UnicodeError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for p in doc.paragraphs:
            if p.style.name.startswith("Heading"):
                lvl = p.style.name.replace("Heading ", "").strip()
                try:
                    n = int(lvl)
                except ValueError:
                    n = 1
                parts.append(f"\n{'#' * min(n, 6)} {p.text}\n")
            else:
                parts.append(p.text)
        for tbl in doc.tables:
            parts.append("")
            for row in tbl.rows:
                parts.append("| " + " | ".join(c.text for c in row.cells) + " |")
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR DOCX: {e}]"


def extract_pdf(path: Path, img_dir: Path) -> tuple[str, int]:
    """Return (text, ocr_count)."""
    text = ""
    ocr_count = 0
    # Try PyMuPDF
    try:
        import fitz
        doc = fitz.open(str(path))
        for i, page in enumerate(doc):
            t = page.get_text()
            if t.strip():
                text += f"\n### Страница {i + 1}\n{t}\n"
            else:
                # No text layer — extract image for OCR
                pix = page.get_pixmap(dpi=150)
                img_path = img_dir / f"{path.stem}_p{i:03d}.png"
                pix.save(str(img_path))
                text += f"\n### Страница {i + 1}\n[изображение — OCR ниже]\n"
        doc.close()
    except ImportError:
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text() or ""
                    text += f"\n### Страница {i + 1}\n{t}\n"
        except ImportError:
            if shutil.which("pdftotext"):
                r = subprocess.run(["pdftotext", "-layout", str(path), "-"],
                                   capture_output=True, text=True)
                text = r.stdout
            else:
                return "[PDF extraction unavailable — install PyMuPDF]", 0

    # OCR for extracted page images
    if img_dir.exists():
        for img in sorted(img_dir.glob("*.png")):
            ocr_text = ocr_image(img)
            if ocr_text.strip():
                text += f"\n[OCR результат]\n{ocr_text}\n"
                ocr_count += 1
            img.unlink()  # clean up

    if not text.strip():
        text = "[PDF пуст — нет текстового слоя и изображений]"
    return text, ocr_count


def extract_pptx(path: Path, img_dir: Path) -> tuple[str, int]:
    parts = []
    ocr_count = 0
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        for n, slide in enumerate(prs.slides, 1):
            parts.append(f"\n## Слайд {n}\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            parts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append("| " + " | ".join(c.text for c in row.cells) + " |")
    except Exception as e:
        parts.append(f"[ERROR PPTX: {e}]")
    return "\n".join(parts), ocr_count


def extract_fb2(path: Path) -> str:
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(str(path))
        root = tree.getroot()
        ns = {"fb": "http://www.gribuser.ru/xml/fictionbook/2.0"}
        parts = []
        for section in root.iter("{http://www.gribuser.ru/xml/fictionbook/2.0}section"):
            title_el = section.find("fb:title", ns)
            if title_el is not None:
                t = "".join(title_el.itertext()).strip()
                if t:
                    parts.append(f"\n## {t}\n")
            for p in section.iter("{http://www.gribuser.ru/xml/fictionbook/2.0}p"):
                t = "".join(p.itertext()).strip()
                if t:
                    parts.append(t)
        return "\n\n".join(parts) if parts else "[FB2: no text extracted]"
    except Exception as e:
        return f"[ERROR FB2: {e}]"


def extract_djvu(path: Path) -> str:
    """Extract text from DJVU using djvused or djvutxt."""
    for tool, args in [
        ("djvutxt", [str(path)]),
        ("djvused", ["-e", "print-txt", str(path)]),
    ]:
        if shutil.which(tool):
            try:
                r = subprocess.run([tool] + args, capture_output=True, text=True, timeout=120)
                if r.stdout.strip():
                    return r.stdout
            except Exception:
                continue
    return f"[DJVU extraction unavailable — install djvulibre: {path.name}]"


def ocr_image(img_path: Path) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(img_path)
        return pytesseract.image_to_string(img, lang="rus+eng")
    except Exception as e:
        return f"[OCR error: {e}]"


def extract_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    clean = re.sub(r"<[^>]+>", " ", raw)
    clean = re.sub(r"\s+", " ", clean).strip()
    return clean


def extract_archive(path: Path, dest: Path) -> list[Path]:
    extracted = []
    try:
        if path.suffix.lower() == ".zip" or path.name.endswith(".zip"):
            with zipfile.ZipFile(path) as zf:
                zf.extractall(dest)
                extracted = [dest / n for n in zf.namelist()]
        elif path.suffix.lower() == ".rar":
            if shutil.which("unrar"):
                subprocess.run(["unrar", "x", str(path), str(dest) + os.sep],
                               check=True, capture_output=True, timeout=120)
                extracted = list(dest.rglob("*"))
            elif shutil.which("7z"):
                subprocess.run(["7z", "x", str(path), f"-o{dest}"],
                               check=True, capture_output=True, timeout=120)
                extracted = list(dest.rglob("*"))
            else:
                log.warning("No unrar/7z — skip %s", path.name)
        elif path.suffix.lower() in (".7z", ".tar", ".tar.gz") or path.name.endswith((".tar.gz", ".tgz")):
            if shutil.which("7z"):
                subprocess.run(["7z", "x", str(path), f"-o{dest}"],
                               check=True, capture_output=True, timeout=120)
                extracted = list(dest.rglob("*"))
            else:
                log.warning("No 7z — skip %s", path.name)
    except Exception as e:
        log.error("Archive error %s: %s", path.name, e)
    return [f for f in extracted if f.is_file()]


# ---------------------------------------------------------------------------
# Classify
# ---------------------------------------------------------------------------

def classify(path: Path) -> str:
    ext = path.suffix.lower()
    if path.name.endswith((".tar.gz", ".tgz")):
        ext = ".tar.gz"
    if ext in ARCHIVE_EXTS:
        return "archive"
    if ext in TEXT_EXTS:
        return "text"
    if ext in DOC_EXTS:
        return "doc"
    if ext in PDF_EXTS:
        return "pdf"
    if ext in PPTX_EXTS:
        return "presentation"
    if ext in BOOK_EXTS:
        return "book"
    if ext == ".djvu":
        return "djvu"
    if ext in HTML_EXTS:
        return "html"
    if ext in IMAGE_EXTS:
        return "image"
    return "skip"


def normalize_for_fingerprint(text: str) -> str:
    normalized = re.sub(r"\s+", " ", text).strip().lower()
    normalized = re.sub(r"[^0-9a-zа-яё]+", "", normalized)
    return normalized


def build_group_slug(rel: Path) -> str:
    parent = rel.parent.as_posix()
    if parent in {"", "."}:
        return ""
    slug = re.sub(r"[^0-9a-zа-яё]+", "-", parent.lower()).strip("-")
    return slug or "source-pack"


def detect_review_flags(content: str, ftype: str, word_count: int) -> list[str]:
    flags = []
    lowered = content.lower()

    if "[error" in lowered or "[unsupported" in lowered or "unavailable" in lowered:
        flags.append("extraction_error")
    if "[ocr" in lowered or "ocr error" in lowered:
        flags.append("ocr_involved")
    if "�" in content:
        flags.append("replacement_chars")
    if ftype in {"pdf", "doc", "presentation", "book", "djvu", "html", "image", "archive"} and word_count < 40:
        flags.append("low_text_yield")
    if re.search(r"[^\w\s]{10,}", content):
        flags.append("noise_suspected")
    return sorted(set(flags))


def make_entry(rel: Path, out_path: Path, output_root: Path, ftype: str, content: str, word_count: int,
               review_flags: list[str], preferred_for_context: bool = True) -> dict:
    prepared_status = "review_needed" if review_flags else "prepared_trusted"
    group_slug = build_group_slug(rel)
    source_rel = rel.as_posix()
    output_rel = out_path.relative_to(output_root).as_posix()
    original_rel = (Path("originals") / rel).as_posix()
    fingerprint = normalize_for_fingerprint(content)
    return {
        "file": output_rel,
        "title": out_path.stem,
        "source_file": source_rel,
        "source_kind": ftype,
        "word_count": word_count,
        "prepared_status": prepared_status,
        "review_flags": review_flags,
        "original_fallback": original_rel,
        "original_fallback_required": prepared_status == "review_needed",
        "preferred_for_context": preferred_for_context and prepared_status == "prepared_trusted",
        "duplicate_of": None,
        "pack_group": group_slug,
        "merged_into": None,
        "_fingerprint": fingerprint,
    }


def build_merged_pack(output: Path, subject: str, group_slug: str, members: list[dict]) -> dict:
    merged_dir = output / "merged-packs"
    merged_dir.mkdir(parents=True, exist_ok=True)
    merged_rel = Path("merged-packs") / f"{group_slug}.md"
    merged_path = output / merged_rel

    chunks = []
    review_flags = set()
    total_words = 0
    for member in members:
        member_path = output / member["file"]
        text = member_path.read_text(encoding="utf-8", errors="replace")
        second_delim = text.find("---", text.find("---", 3) + 3)
        body = text[second_delim + 3:].strip() if second_delim > 0 else text.strip()
        chunks.append(
            f"## {member['title']}\n\n"
            f"> Source file: `{member['source_file']}`\n\n"
            f"{body}\n"
        )
        review_flags.update(member["review_flags"])
        total_words += member["word_count"]

    metadata = {
        "source_kind": "merged_pack",
        "prepared_status": "review_needed" if review_flags else "prepared_trusted",
        "review_flags": sorted(review_flags),
        "original_fallback": "",
        "original_fallback_required": bool(review_flags),
        "pack_group": group_slug,
        "merged_from": [member["file"] for member in members],
    }
    merged_content = "\n---\n".join(chunks)
    merged_md = build_md(merged_rel, merged_content, subject, merged_path, total_words, metadata)
    merged_path.write_text(merged_md, encoding="utf-8")

    return {
        "file": merged_rel.as_posix(),
        "title": merged_path.stem,
        "source_file": "",
        "source_kind": "merged_pack",
        "word_count": total_words,
        "prepared_status": metadata["prepared_status"],
        "review_flags": sorted(review_flags),
        "original_fallback": "",
        "original_fallback_required": bool(review_flags),
        "preferred_for_context": not review_flags,
        "duplicate_of": None,
        "pack_group": group_slug,
        "merged_into": None,
        "merged_from": [member["file"] for member in members],
        "_fingerprint": normalize_for_fingerprint(merged_content),
    }


def finalize_entries(output: Path, subject: str, stats: dict) -> None:
    fingerprint_first = {}
    grouped = defaultdict(list)

    for entry in sorted(stats["entries"], key=lambda item: item["file"]):
        fingerprint = entry.get("_fingerprint") or ""
        if fingerprint and len(fingerprint) >= 80:
            if fingerprint in fingerprint_first:
                entry["duplicate_of"] = fingerprint_first[fingerprint]
                entry["preferred_for_context"] = False
            else:
                fingerprint_first[fingerprint] = entry["file"]

        if entry["duplicate_of"] is None and entry["pack_group"] and entry["word_count"] <= 220:
            grouped[entry["pack_group"]].append(entry)

    merged_entries = []
    for group_slug, members in grouped.items():
        if len(members) < 2:
            continue
        total_words = sum(member["word_count"] for member in members)
        if total_words > 1200:
            continue
        merged_entry = build_merged_pack(output, subject, group_slug, members)
        for member in members:
            member["merged_into"] = merged_entry["file"]
            member["preferred_for_context"] = False
        merged_entries.append(merged_entry)
        stats["output_files"].append(merged_entry["file"])

    stats["entries"].extend(merged_entries)


# ---------------------------------------------------------------------------
# Main processing
# ---------------------------------------------------------------------------

def process_dir(source: Path, output: Path, subject: str) -> dict:
    stats = {
        "files_scanned": 0, "archives": 0, "text": 0, "doc": 0,
        "pdf": 0, "pptx": 0, "book": 0, "djvu": 0, "html": 0,
        "image": 0, "skipped": 0, "errors": 0, "output_files": [],
        "total_words": 0, "entries": [],
    }
    img_dir = output / "_ocr_tmp"
    img_dir.mkdir(parents=True, exist_ok=True)

    for file_path in sorted(source.rglob("*")):
        if not file_path.is_file():
            continue

        # Skip hidden/system dirs
        skip = False
        for part in file_path.parts:
            if part in SKIP_DIRS:
                skip = True
                break
        if skip:
            continue

        stats["files_scanned"] += 1
        ftype = classify(file_path)
        rel = file_path.relative_to(source)

        if ftype == "skip" or file_path.suffix.lower() in SKIP_EXTS:
            stats["skipped"] += 1
            continue

        log.info("[%s] %s", ftype.upper(), rel)

        content = ""
        try:
            if ftype == "text":
                content = read_text(file_path)
                stats["text"] += 1

            elif ftype == "doc":
                content = extract_docx(file_path)
                stats["doc"] += 1

            elif ftype == "pdf":
                content, ocr_n = extract_pdf(file_path, img_dir)
                stats["pdf"] += 1

            elif ftype == "presentation":
                content, _ = extract_pptx(file_path, img_dir)
                stats["pptx"] += 1

            elif ftype == "book":
                ext = file_path.suffix.lower()
                if ext == ".fb2":
                    content = extract_fb2(file_path)
                else:
                    content = f"[Unsupported book format: {ext}]"
                stats["book"] += 1

            elif ftype == "djvu":
                content = extract_djvu(file_path)
                stats["djvu"] += 1

            elif ftype == "html":
                content = extract_html(file_path)
                stats["html"] += 1

            elif ftype == "image":
                content = ocr_image(file_path)
                stats["image"] += 1

            elif ftype == "archive":
                tmp = tempfile.mkdtemp(prefix="study-")
                files = extract_archive(file_path, Path(tmp))
                stats["archives"] += 1

                # Save extracted originals to output/extracted/
                extracted_dir = output / "extracted" / rel.stem
                extracted_dir.mkdir(parents=True, exist_ok=True)
                for ef in files:
                    if ef.is_file():
                        shutil.copy2(ef, extracted_dir / ef.name)

                # Process extracted files inline
                for ef in files:
                    ef_type = classify(ef)
                    if ef_type == "skip":
                        continue
                    ec = ""
                    if ef_type == "text":
                        ec = read_text(ef)
                    elif ef_type == "doc":
                        ec = extract_docx(ef)
                    elif ef_type == "pdf":
                        ec, _ = extract_pdf(ef, img_dir)
                    elif ef_type == "presentation":
                        ec, _ = extract_pptx(ef, img_dir)
                    elif ef_type == "book":
                        ec = extract_fb2(ef) if ef.suffix.lower() == ".fb2" else f"[{ef.suffix}]"
                    elif ef_type == "djvu":
                        ec = extract_djvu(ef)
                    elif ef_type == "html":
                        ec = extract_html(ef)
                    elif ef_type == "image":
                        ec = ocr_image(ef)
                    if ec.strip():
                        content += f"\n---\n### Файл: {ef.name}\n{ec}\n"
                shutil.rmtree(tmp, ignore_errors=True)

        except Exception as e:
            log.error("Error %s: %s", rel, e)
            stats["errors"] += 1
            content = f"[ERROR: {e}]"

        if not content.strip():
            continue

        # Save original file copy
        orig_dir = output / "originals" / rel.parent
        orig_dir.mkdir(parents=True, exist_ok=True)
        shutil.copy2(file_path, orig_dir / file_path.name)

        # Build output Markdown
        word_count = len(content.split())
        stats["total_words"] += word_count
        out_path = build_out_path(output, rel, file_path)
        review_flags = detect_review_flags(content, ftype, word_count)
        entry = make_entry(rel, out_path, output, ftype, content, word_count, review_flags)
        md = build_md(rel, content, subject, file_path, word_count, entry)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(md, encoding="utf-8")
        stats["output_files"].append(str(out_path.relative_to(output)))
        stats["entries"].append(entry)

    # Clean up OCR temp
    shutil.rmtree(img_dir, ignore_errors=True)
    finalize_entries(output, subject, stats)
    return stats


def build_md(rel: Path, content: str, subject: str, orig: Path, wc: int, metadata: dict | None = None) -> str:
    metadata = metadata or {}
    lines = [
        "---",
        f"title: \"{rel.stem.replace('_', ' ').strip()}\"",
        f"subject: \"{subject}\"",
        f"source_file: \"{rel.as_posix()}\"",
        f"source_kind: \"{metadata.get('source_kind', 'unknown')}\"",
        f"processed_at: \"{datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%SZ')}\"",
        f"original_hash: \"{hashlib.sha256(content.encode()).hexdigest()[:16]}\"",
        f"word_count: {wc}",
        f"prepared_status: \"{metadata.get('prepared_status', 'prepared_trusted')}\"",
        "review_flags:",
        "---",
        "",
        f"# {rel.stem.replace('_', ' ').strip()}",
        "",
        f"> Source: `{rel.as_posix()}`",
        "",
        "---",
        "",
        content,
        "",
    ]
    review_flags = metadata.get("review_flags", [])
    if review_flags:
        review_idx = lines.index("review_flags:") + 1
        for flag in review_flags:
            lines.insert(review_idx, f"  - \"{flag}\"")
            review_idx += 1
    else:
        review_idx = lines.index("review_flags:") + 1
        lines.insert(review_idx, "  - \"none\"")

    insert_at = lines.index("---", 1)
    optional_lines = [
        f"original_fallback: \"{metadata.get('original_fallback', '')}\"",
        f"original_fallback_required: {str(metadata.get('original_fallback_required', False)).lower()}",
    ]
    if metadata.get("pack_group"):
        optional_lines.append(f"pack_group: \"{metadata['pack_group']}\"")
    if metadata.get("duplicate_of"):
        optional_lines.append(f"duplicate_of: \"{metadata['duplicate_of']}\"")
    if metadata.get("merged_into"):
        optional_lines.append(f"merged_into: \"{metadata['merged_into']}\"")
    if metadata.get("merged_from"):
        optional_lines.append("merged_from:")
        optional_lines.extend(f"  - \"{item}\"" for item in metadata["merged_from"])
    for offset, line in enumerate(optional_lines):
        lines.insert(insert_at + offset, line)
    return "\n".join(lines)


def build_out_path(output: Path, rel: Path, orig: Path) -> Path:
    # Preserve directory structure, change extension to .md
    stem = rel.stem
    return output / rel.parent / (stem + ".md")


def write_index(output: Path, stats: dict, subject: str):
    entries = sorted(stats["entries"], key=lambda item: item["file"])
    preferred_context_files = [
        entry["file"] for entry in entries
        if entry["preferred_for_context"] and not entry["duplicate_of"]
    ]
    review_before_use_files = [
        entry["file"] for entry in entries
        if entry["prepared_status"] == "review_needed"
    ]
    duplicate_files = [
        {"file": entry["file"], "duplicate_of": entry["duplicate_of"]}
        for entry in entries if entry["duplicate_of"]
    ]
    merged_packs = [
        {"file": entry["file"], "merged_from": entry.get("merged_from", [])}
        for entry in entries if entry["source_kind"] == "merged_pack"
    ]

    index = {
        "subject": subject,
        "generated_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "files_scanned": stats["files_scanned"],
        "output_files": len(stats["output_files"]),
        "total_words": stats["total_words"],
        "statistics": {k: v for k, v in stats.items() if k not in {"output_files", "entries"}},
        "ingestion_workflow": {
            "launch_mode": "agent_launched_orchestrated_step",
            "verification_required": True,
            "originals_retained": True,
            "duplicate_strategy": "mark_duplicate_and_reduce_preferred_context",
            "consolidation_strategy": "safe_group_merge_for_small_related_fragments",
        },
        "preferred_context_files": preferred_context_files,
        "review_before_use_files": review_before_use_files,
        "duplicate_files": duplicate_files,
        "merged_packs": merged_packs,
        "entries": [],
    }
    for entry in entries:
        md_file = output / entry["file"]
        text = md_file.read_text(encoding="utf-8", errors="replace")
        idx = text.find("---", text.find("---", 3) + 3)
        preview = ""
        if idx > 0:
            preview = text[idx + 3:].strip()[:300]
        clean_entry = {k: v for k, v in entry.items() if not k.startswith("_")}
        clean_entry["preview"] = preview.replace("\n", " ")
        index["entries"].append(clean_entry)
    (output / "index.json").write_text(
        json.dumps(index, indent=2, ensure_ascii=False), encoding="utf-8")


def write_readme(output: Path, stats: dict, subject: str, source: Path):
    trusted = [entry for entry in stats["entries"] if entry["prepared_status"] == "prepared_trusted"]
    review_needed = [entry for entry in stats["entries"] if entry["prepared_status"] == "review_needed"]
    duplicates = [entry for entry in stats["entries"] if entry["duplicate_of"]]
    merged = [entry for entry in stats["entries"] if entry["source_kind"] == "merged_pack"]

    lines = [
        f"# {subject} — Study Materials (processed)",
        "",
        f"Source: `{source}`",
        f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}",
        "",
        "## Summary",
        "",
        f"- Files scanned: {stats['files_scanned']}",
        f"- Archives extracted: {stats['archives']}",
        f"- Text files: {stats['text']}",
        f"- DOCX: {stats['doc']}",
        f"- PDF: {stats['pdf']}",
        f"- PPTX: {stats['pptx']}",
        f"- Books (FB2): {stats['book']}",
        f"- DJVU: {stats['djvu']}",
        f"- HTML: {stats['html']}",
        f"- Images OCR'd: {stats['image']}",
        f"- Output Markdown files: {len(stats['output_files'])}",
        f"- Total words: {stats['total_words']}",
        f"- Trusted prepared files: {len(trusted)}",
        f"- Review-needed files: {len(review_needed)}",
        f"- Duplicate-marked files: {len(duplicates)}",
        f"- Merged packs: {len(merged)}",
        "",
        "## Output structure",
        "",
        "Each file preserves the original relative directory structure.",
        "All files have YAML front matter with subject, source, hash, word count, prepared status, review flags, and original fallback path.",
        "Original source files are preserved under `originals/` for fallback when conversion is weak or suspicious.",
        "",
        "## Agent-launched ingestion workflow",
        "",
        "Use this prep step as an agent/workflow action:",
        "1. Inspect the incoming material set and decide whether preparation is recommended.",
        "2. Run `scripts/study-materials-prep.py`.",
        "3. Review `index.json`, the generated Markdown, and the review-needed list before trusting the pack.",
        "4. Prefer trusted prepared files for downstream skills, and keep originals available when conversion is weak.",
        "5. Use merged packs for overview context when they are present, but fall back to source members or originals for critical claims.",
        "",
        "## Verification and fallback",
        "",
        "- `prepared_trusted` means the extracted text is suitable as the preferred context shape.",
        "- `review_needed` means the launching agent/workflow must review the Markdown and may need to consult `originals/` before using it for precise claims.",
        "- `original_fallback_required: true` means the original file should remain prominent in the academic context packet.",
        "- `duplicate_of` means the file is indexed but should not usually be preferred for downstream context loading.",
        "- merged packs under `merged-packs/` reduce fragmentation for related small files.",
        "",
        "## Usage",
        "",
        "- Load `.md` files into Claude / Codex / Gemini as context",
        "- Use them as prepared source packs for `homework-management` and `case-analyzer`",
        "- Use extracted lecture text as input to `lecture-transcript` when raw lecture files needed OCR or archive extraction",
        "- Upload to ChatGPT as Knowledge files",
        "- Import into NotebookLM as sources",
        "- Use `index.json` for RAG pipeline ingestion",
        "",
    ]
    if review_needed:
        lines.extend(["## Review needed before trusting for precise academic claims", ""])
        for entry in review_needed[:15]:
            flags = ", ".join(entry["review_flags"]) or "manual_review"
            lines.append(
                f"- `{entry['file']}` -> review flags: {flags}; fallback: `{entry['original_fallback']}`"
            )
        lines.append("")
    if merged:
        lines.extend(["## Merged packs", ""])
        for entry in merged[:15]:
            lines.append(f"- `{entry['file']}` from {', '.join(entry.get('merged_from', []))}")
        lines.append("")
    if stats["errors"]:
        lines.extend(["## Errors", "", f"- {stats['errors']} file(s) had extraction errors", ""])
    (output / "README.md").write_text("\n".join(lines), encoding="utf-8")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Prepare study materials for RAG indexing. "
                    "Recursively scans source, extracts all content, outputs structured Markdown.",
        epilog="""
Examples:
  python study-materials-prep.py --source "C:\\path\\to\\HR-менеджмент"
  python study-materials-prep.py --source "C:\\path\\to\\module3\\HR-менеджмент" --output "C:\\output"
        """,
    )
    parser.add_argument("--source", required=True,
                        help="Source directory with study materials")
    parser.add_argument("--output", default="",
                        help="Output directory (default: ./study-output/<subject>/)")
    parser.add_argument("--install-deps", action="store_true",
                        help="Auto-install missing dependencies before processing")
    args = parser.parse_args()

    source = Path(args.source).expanduser().resolve()
    if not source.is_dir():
        log.error("Source directory not found: %s", source)
        sys.exit(1)

    subject = source.name
    if args.output:
        output = Path(args.output).expanduser().resolve()
    else:
        output = Path.cwd() / "study-output" / subject

    warnings = check_deps()
    for w in warnings:
        log.warning(w)

    # Auto-install dependencies if requested or if critical deps missing
    if args.install_deps or any("MISSING:" in w for w in warnings[:2]):
        log.info("Installing missing dependencies...")
        script_dir = Path(__file__).parent
        install_script = script_dir / "install-deps.py"
        if install_script.exists():
            subprocess.run([sys.executable, str(install_script)], timeout=600)
            # Re-check after install
            warnings = check_deps()
            for w in warnings:
                log.warning(w)
        else:
            log.warning("install-deps.py not found — run manually: python scripts/install-deps.py")

    output.mkdir(parents=True, exist_ok=True)

    log.info("Subject: %s", subject)
    log.info("Source: %s", source)
    log.info("Output: %s", output)

    stats = process_dir(source, output, subject)
    write_index(output, stats, subject)
    write_readme(output, stats, subject, source)

    log.info("=" * 60)
    log.info("DONE — %s", subject)
    log.info("=" * 60)
    log.info("Files scanned:     %d", stats["files_scanned"])
    log.info("Archives:          %d", stats["archives"])
    log.info("Text files:        %d", stats["text"])
    log.info("DOCX:              %d", stats["doc"])
    log.info("PDF:               %d", stats["pdf"])
    log.info("PPTX:              %d", stats["pptx"])
    log.info("Books (FB2):       %d", stats["book"])
    log.info("DJVU:              %d", stats["djvu"])
    log.info("HTML:              %d", stats["html"])
    log.info("Images OCR'd:      %d", stats["image"])
    log.info("Output files:      %d", len(stats["output_files"]))
    log.info("Total words:       %d", stats["total_words"])
    log.info("Errors:            %d", stats["errors"])
    log.info("Output directory:  %s", output)


if __name__ == "__main__":
    main()
